"""Word、PDF、Excel 和 PowerPoint 报告导出辅助函数。"""

from __future__ import annotations

import base64
import os
from io import BytesIO
from tempfile import NamedTemporaryFile
from typing import Any, Dict, List, Optional

import pandas as pd
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn

from app.services.file_preview import DATA_CACHE
from app.services.report_builder import build_report


def _safe_cell_to_string(value: Any) -> str:
    """计算 Excel 列宽时避免 pandas NA 值触发错误。"""
    try:
        if pd.isna(value):
            return ""
    except (TypeError, ValueError):
        pass
    return str(value)


def _decode_chart_images(charts: list[dict]) -> list[tuple[str, bytes, str]]:
    """把前端 ECharts data URL 解码为原始图片字节。"""
    images: list[tuple[str, bytes, str]] = []
    for chart in charts or []:
        data_url = chart.get("data_url") or chart.get("dataUrl") or chart.get("image")
        if not data_url:
            continue
        title = chart.get("title") or "Chart"
        ext = "png"
        payload = data_url
        if data_url.startswith("data:"):
            header, payload = data_url.split(",", 1)
            if "image/" in header:
                ext = header.split("image/")[1].split(";")[0].strip()
        try:
            raw = base64.b64decode(payload)
        except Exception:
            continue
        images.append((title, raw, ext))
    return images


def _get_report(saved_name: str) -> dict:
    """基于缓存 DataFrame 重新生成最新报告。"""
    dataframe = DATA_CACHE.get(saved_name)
    if dataframe is None:
        raise ValueError("Session expired or file not found. Please re-upload.")
    return build_report(dataframe)


def _set_cell_shading(cell, color: str):
    shading_elm = cell._element.get_or_add_tcPr()
    shading = shading_elm.makeelement(qn("w:shd"), {
        qn("w:fill"): color,
        qn("w:val"): "clear",
    })
    shading_elm.append(shading)


def export_report_docx(
    saved_name: str,
    original_filename: str,
    charts: Optional[List[Dict[str, Any]]] = None,
) -> BytesIO:
    """创建包含关键表格和可选图表图片的 Word 报告。"""
    report = _get_report(saved_name)
    doc = Document()
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Arial"
    font.size = Pt(10)

    title = doc.add_heading(f"DataFlowBI Analysis Report", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph(f"Source file: {original_filename}")
    doc.add_paragraph("")

    sections = [
        ("Missing Statistics", "missing", None),
        ("Numeric Summary", "numeric_summary", ["count", "mean", "std", "min", "max"]),
        ("Sample Data (First 5 Rows)", "sample_rows", None),
    ]

    for section_title, report_key, col_subset in sections:
        data = report.get(report_key)
        if not data:
            continue

        doc.add_heading(section_title, level=1)

        if report_key == "missing":
            table = doc.add_table(rows=1, cols=3)
            table.style = "Light Grid Accent 1"
            hdr = table.rows[0].cells
            hdr[0].text = "Field"
            hdr[1].text = "Missing Count"
            hdr[2].text = "Missing Rate"
            for field in sorted(data.keys()):
                row_cells = table.add_row().cells
                row_cells[0].text = str(field)
                row_cells[1].text = str(data[field])
                mr = report.get("missing_rate", {}).get(field, 0)
                row_cells[2].text = f"{mr:.2%}" if isinstance(mr, (int, float)) else str(mr)

        elif report_key == "numeric_summary":
            if isinstance(data, dict) and data:
                cols = col_subset or ["count", "mean", "std", "min", "max"]
                table = doc.add_table(rows=1, cols=1 + len(cols))
                table.style = "Light Grid Accent 1"
                hdr = table.rows[0].cells
                hdr[0].text = "Field"
                for i, c in enumerate(cols):
                    hdr[i + 1].text = c.capitalize()
                for field, stats in sorted(data.items()):
                    row_cells = table.add_row().cells
                    row_cells[0].text = str(field)
                    for i, c in enumerate(cols):
                        val = stats.get(c, "")
                        row_cells[i + 1].text = f"{val:.4f}" if isinstance(val, float) else str(val)

        elif report_key == "sample_rows":
            if isinstance(data, list) and data:
                cols = list(data[0].keys())
                table = doc.add_table(rows=1, cols=len(cols))
                table.style = "Light Grid Accent 1"
                hdr = table.rows[0].cells
                for i, c in enumerate(cols):
                    hdr[i].text = str(c)
                for row in data:
                    row_cells = table.add_row().cells
                    for i, c in enumerate(cols):
                        row_cells[i].text = str(row.get(c, ""))

        doc.add_paragraph("")

    if charts:
        doc.add_heading("Charts", level=1)
        for title, image_bytes, _ in _decode_chart_images(charts):
            doc.add_paragraph(title)
            doc.add_picture(BytesIO(image_bytes), width=Inches(5.8))
            doc.add_paragraph("")

    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def export_report_pdf(
    saved_name: str,
    original_filename: str,
    charts: Optional[List[Dict[str, Any]]] = None,
) -> BytesIO:
    """创建包含摘要表格和图表的紧凑 PDF 报告。"""
    report = _get_report(saved_name)
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 12, "DataFlowBI Analysis Report", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.set_font("Helvetica", "", 10)
    pdf.cell(0, 8, f"Source file: {original_filename}", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(4)

    def write_table(headers, rows, col_widths=None):
        if not col_widths:
            col_widths = [pdf.w / (len(headers) + 1)] * len(headers)
        pdf.set_font("Helvetica", "B", 9)
        for i, h in enumerate(headers):
            pdf.cell(col_widths[i], 8, str(h), border=1)
        pdf.ln()
        pdf.set_font("Helvetica", "", 8)
        for row in rows:
            for i, cell in enumerate(row):
                pdf.cell(col_widths[i], 7, str(cell)[:30], border=1)
            pdf.ln()

    missing = report.get("missing", {})
    if missing:
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 10, "Missing Statistics", new_x="LMARGIN", new_y="NEXT")
        mr = report.get("missing_rate", {})
        headers = ["Field", "Missing Count", "Missing Rate"]
        rows = [[f, missing[f], f"{mr.get(f, 0):.2%}"] for f in sorted(missing)]
        write_table(headers, rows)

    ns = report.get("numeric_summary", {})
    if ns:
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 10, "Numeric Summary", new_x="LMARGIN", new_y="NEXT")
        headers = ["Field", "Count", "Mean", "Std", "Min", "Max"]
        rows = []
        for f, s in sorted(ns.items()):
            rows.append([f, s.get("count", ""), f"{s.get('mean', ''):.4f}" if isinstance(s.get("mean"), float) else str(s.get("mean", "")),
                         f"{s.get('std', ''):.4f}" if isinstance(s.get("std"), float) else str(s.get("std", "")),
                         f"{s.get('min', ''):.4f}" if isinstance(s.get("min"), float) else str(s.get("min", "")),
                         f"{s.get('max', ''):.4f}" if isinstance(s.get("max"), float) else str(s.get("max", ""))])
        write_table(headers, rows)

    sample = report.get("sample_rows", [])
    if sample:
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 10, "Sample Data (First 5 Rows)", new_x="LMARGIN", new_y="NEXT")
        headers = list(sample[0].keys())
        rows = [[str(r.get(c, ""))[:20] for c in headers] for r in sample]
        write_table(headers, rows, col_widths=[pdf.w / (len(headers) + 1)] * len(headers))

    if charts:
        pdf.add_page()
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 10, "Charts", new_x="LMARGIN", new_y="NEXT")
        for title, image_bytes, ext in _decode_chart_images(charts):
            pdf.set_font("Helvetica", "B", 10)
            pdf.cell(0, 8, title, new_x="LMARGIN", new_y="NEXT")
            tmp_path = None
            try:
                with NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmp:
                    tmp.write(image_bytes)
                    tmp.flush()
                    tmp_path = tmp.name
                pdf.image(tmp_path, w=pdf.w - 20)
            finally:
                if tmp_path and os.path.exists(tmp_path):
                    os.remove(tmp_path)
            pdf.ln(4)

    buf = BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return buf


def export_report_excel(saved_name: str, original_filename: str) -> BytesIO:
    """根据报告分区创建多工作表 Excel 文件。"""
    report = _get_report(saved_name)
    buf = BytesIO()

    sheet_map: list[tuple[str, Any]] = [
        ("Missing Stats", report.get("missing")),
        ("Numeric Summary", report.get("numeric_summary")),
        ("Sample Data", report.get("sample_rows")),
        ("Frequency", report.get("frequencies")),
        ("Pareto", report.get("pareto")),
        ("Boxplot", report.get("boxplot")),
        ("Correlation", report.get("correlation")),
        ("Group Stats", report.get("group_stats")),
        ("Binning", report.get("binning")),
        ("Violin", report.get("violin")),
        ("Scatter Matrix", report.get("scatter_matrix")),
        ("Missing Heatmap", report.get("missing_heatmap")),
        ("Time Series", report.get("timeseries")),
        ("Outliers", report.get("outliers")),
    ]

    with pd.ExcelWriter(buf, engine="xlsxwriter") as writer:
        workbook = writer.book
        header_fmt = workbook.add_format({"bold": True, "bg_color": "#f3f4f6", "border": 1})
        for sheet_name, data in sheet_map:
            if not data:
                continue

            safe_name = sheet_name[:31]
            df = None
            if isinstance(data, list):
                df = pd.DataFrame(data)
            elif isinstance(data, dict):
                try:
                    df = pd.DataFrame(data).transpose().reset_index()
                    df.columns = ["Field"] + [str(c) for c in df.columns[1:]]
                except (ValueError, TypeError):
                    rows = []
                    for k, v in data.items():
                        if isinstance(v, dict):
                            row = {"Field": k}
                            row.update(v)
                            rows.append(row)
                        else:
                            rows.append({"Field": k, "Value": v})
                    if rows:
                        df = pd.DataFrame(rows)

            if df is None or df.empty:
                continue

            df.to_excel(writer, sheet_name=safe_name, index=False)
            worksheet = writer.sheets[safe_name]
            for col_num, value in enumerate(df.columns.values):
                worksheet.write(0, col_num, value, header_fmt)
                series = df.iloc[:, col_num].map(_safe_cell_to_string)
                max_len = max(series.map(len).max(), len(str(value))) + 2
                worksheet.set_column(col_num, col_num, min(max_len, 40))

    buf.seek(0)
    return buf


def export_report_pptx(
    saved_name: str,
    original_filename: str,
    charts: Optional[List[Dict[str, Any]]] = None,
) -> BytesIO:
    """创建包含摘要页和图表页的 PowerPoint 文件。"""
    from pptx import Presentation
    from pptx.util import Inches

    report = _get_report(saved_name)
    images = _decode_chart_images(charts or [])
    prs = Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "DataFlowBI Analysis Report"
    title_slide.placeholders[1].text = f"Source file: {original_filename}"

    summary_slide = prs.slides.add_slide(prs.slide_layouts[5])
    summary_slide.shapes.title.text = "Summary"
    body = summary_slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(4.5)).text_frame
    field_count = len(report.get("dtypes", {}))
    body.text = f"Fields: {field_count}"

    if images:
        for title, image_bytes, _ in images:
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
            slide.shapes.add_picture(BytesIO(image_bytes), Inches(0.6), Inches(1.6), width=Inches(8.8))

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
